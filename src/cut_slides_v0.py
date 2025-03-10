
from pptx import Presentation

# Load the original PowerPoint file
input_pptx = "input_presentation.pptx"  # Replace with your filename
presentation = Presentation(input_pptx)

# Loop through each slide and save it as a separate PowerPoint file
for i, slide in enumerate(presentation.slides):
    new_presentation = Presentation()
    
    # Copy slide layout from the original presentation
    slide_layout = new_presentation.slide_layouts[5]  # Use a blank layout
    new_slide = new_presentation.slides.add_slide(slide_layout)

    # Copy shapes (text, images, etc.) from the original slide
    for shape in slide.shapes:
        if shape.has_text_frame:
            new_shape = new_slide.shapes.add_textbox(
                left=shape.left, top=shape.top, width=shape.width, height=shape.height
            )
            new_shape.text = shape.text

    # Save the individual slide as a new PowerPoint file
    output_pptx = f"slide_{i+1}.pptx"
    new_presentation.save(output_pptx)

print("Slides have been saved as separate PowerPoint files.")